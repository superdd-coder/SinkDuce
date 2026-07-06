from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from src.parsers.base import DocumentParser, ImageInfo, ParsedDocument
import uuid


class PptxParser(DocumentParser):
    def parse(self, path: Path) -> ParsedDocument:
        prs = Presentation(str(path))
        slides_text: list[str] = []
        all_images: list[ImageInfo] = []

        for i, slide in enumerate(prs.slides):
            slide_images: list[ImageInfo] = []
            texts: list[str] = []

            for shape in slide.shapes:
                if shape.has_text_frame:
                    texts.append(shape.text_frame.text)
                elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    # Extract image
                    try:
                        image = shape.image
                        img_id = uuid.uuid4().hex
                        slide_images.append(ImageInfo(
                            image_id=img_id,
                            slide_number=i + 1,
                            # Convert EMU to approximate px (1 inch = 914400 EMU, ~96 DPI → 1 px ≈ 9525 EMU)
                            bbox=(
                                shape.left / 9525, shape.top / 9525,
                                (shape.left + shape.width) / 9525, (shape.top + shape.height) / 9525,
                            ),
                            image_bytes=image.blob,
                            image_format=image.content_type.split("/")[-1] if image.content_type else "png",
                            alt_text=getattr(shape, "alt_text", "") or "",
                        ))
                    except Exception:
                        pass

            # Build slide text
            slide_content_parts: list[str] = []
            if texts:
                slide_content_parts.append(f"## Slide {i + 1}\n" + "\n".join(texts))

            # Append :::image blocks at the end of each slide
            for img in slide_images:
                block = (
                    f":::image\n"
                    f"image_id: {img.image_id}\n"
                    f"file_id: \n"  # filled by upload_handler
                    f"description: \n"  # filled by describe step
                    f":::"
                )
                slide_content_parts.append(block)

            if slide_content_parts:
                slides_text.append("\n\n".join(slide_content_parts))
            all_images.extend(slide_images)

        # Build content with position_map tracking slide boundaries
        position_map: list[dict] = []
        parts: list[str] = []
        offset = 0
        for i, slide_text in enumerate(slides_text):
            position_map.append({
                "char_offset": offset,
                "label": f"Slide {i + 1}",
                "type": "slide",
                "slide_number": i + 1,
            })
            parts.append(slide_text)
            offset += len(slide_text) + 2

        return ParsedDocument(
            content="\n\n".join(parts),
            metadata={"slides": len(slides_text)},
            source_path=str(path),
            file_type="pptx",
            position_map=position_map,
            # Images NOT filtered here — process_document_images in upload_handler
            # handles filtering + block cleanup atomically.
            images=all_images,
        )
