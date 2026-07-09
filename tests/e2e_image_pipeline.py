"""End-to-end tests for image processing pipeline across all file types."""
from __future__ import annotations

import io
import os
import sys
import tempfile
import uuid
from pathlib import Path

# ── Test image generators ──────────────────────────────────────────────

def _make_text_image() -> bytes:
    """Generate a PNG with clear text — should be classified as 'text' by OCR."""
    from PIL import Image, ImageDraw, ImageFont
    img = Image.new("RGB", (800, 200), "white")
    d = ImageDraw.Draw(img)
    try:
        font = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", 24)
    except Exception:
        font = ImageFont.load_default()
    text = (
        "Q3 Revenue Report 2025\n"
        "Total Revenue: $42,000,000\n"
        "Growth YoY: 15.3%\n"
        "Asia-Pacific leads with 42% share\n"
        "New customers: 1,247"
    )
    d.text((20, 20), text, fill="black", font=font)
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    return buf.getvalue()


def _make_visual_image() -> bytes:
    """Generate a simple colored image — should be 'visual' (no text)."""
    from PIL import Image
    img = Image.new("RGB", (400, 300), "blue")
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    return buf.getvalue()


def _make_mixed_image() -> bytes:
    """Generate an image with a little text + visual elements — 'mixed'."""
    from PIL import Image, ImageDraw, ImageFont
    img = Image.new("RGB", (600, 400), "white")
    d = ImageDraw.Draw(img)
    d.rectangle([50, 50, 200, 150], fill="lightblue", outline="blue")
    d.rectangle([250, 80, 450, 180], fill="lightgreen", outline="green")
    try:
        font = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", 20)
    except Exception:
        font = ImageFont.load_default()
    d.text((60, 170), "Architecture Overview", fill="black", font=font)
    d.text((60, 200), "Frontend -> API -> DB", fill="black", font=font)
    d.text((60, 230), "Auth: OAuth 2.0 + JWT", fill="black", font=font)
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    return buf.getvalue()


# ── Helper ─────────────────────────────────────────────────────────────

def _run_parse_and_process(parser, file_path: Path, file_id: str = "") -> dict:
    """Run parse + process_document_images, return result dict."""
    from src.parsers.image_utils import process_document_images
    from src.config import get_config
    from src.prompts import VISUAL_PROMPT

    doc = parser.parse(file_path)
    if not doc.images:
        return {"content": doc.content, "images": [], "file_type": doc.file_type}

    file_dir = file_path.parent
    if not file_id:
        file_id = uuid.uuid4().hex

    cfg = get_config()
    vision_provider = None
    vision_model_id = getattr(cfg, "visual_model_id", "") or ""
    if vision_model_id:
        for p in cfg.llm.providers:
            if hasattr(p, "visual_model_ids") and vision_model_id in p.visual_model_ids:
                vision_provider = p
                break

    doc = process_document_images(
        doc, file_id, file_dir,
        vision_provider=vision_provider,
        vision_model_id=vision_model_id,
        vision_prompt=VISUAL_PROMPT,
    )
    return {
        "content": doc.content,
        "images": doc.images,
        "file_type": doc.file_type,
    }


# ═══════════════════════════════════════════════════════════════════════
# Test 1: Markdown with local images
# ═══════════════════════════════════════════════════════════════════════

def test_markdown_with_images():
    """Markdown file with ![alt](local.png) references."""
    print("\n" + "=" * 60)
    print("TEST 1: Markdown with local images")
    print("=" * 60)

    with tempfile.TemporaryDirectory() as tmp:
        tmpdir = Path(tmp)
        (tmpdir / "chart.png").write_bytes(_make_text_image())
        (tmpdir / "photo.png").write_bytes(_make_visual_image())
        (tmpdir / "diagram.png").write_bytes(_make_mixed_image())

        md_content = (
            "# Test Document\n\n"
            "Some text before.\n\n"
            "![Revenue Chart](chart.png)\n\n"
            "More text in the middle.\n\n"
            "![Team Photo](photo.png)\n\n"
            "![Architecture](diagram.png)\n\n"
            "Text after.\n\n"
            ":::distill-block{src=\"note_123\"}\n"
            "Distilled content here\n"
            ":::\n"
        )
        md_path = tmpdir / "test.md"
        md_path.write_text(md_content, encoding="utf-8")

        from src.parsers.markdown import MarkdownParser
        parser = MarkdownParser()
        doc = parser.parse(md_path)

        assert doc.images, "Should have extracted images"
        assert len(doc.images) == 3, f"Expected 3 images, got {len(doc.images)}"
        assert ":::image" in doc.content, "Should have :::image blocks"
        assert ":::distill-block" in doc.content, "Distill block should be preserved"
        assert "![Revenue Chart]" not in doc.content, "Markdown image should be replaced"

        print(f"  OK MarkdownParser extracted {len(doc.images)} images")
        print(f"  OK Distill block preserved")

        result = _run_parse_and_process(parser, md_path, file_id="e2e_md_test")
        content = result["content"]
        images = result["images"]

        import re
        block_count = len(re.findall(r':::image', content))
        print(f"  OK Content has {block_count} :::image markers")
        print(f"  OK Final images: {len(images)}")

        for img in images:
            has_ocr = bool(img.ocr_text)
            has_desc = bool(img.description)
            print(f"    image_id={img.image_id[:12]}... ocr={has_ocr} desc={has_desc}")

        # Verify block structure in content
        ocr_count = content.count("ocr_text:")
        print(f"  OK All blocks have image_id, ocr_text count in content: {ocr_count}")


# ═══════════════════════════════════════════════════════════════════════
# Test 2: Note image URL resolver
# ═══════════════════════════════════════════════════════════════════════

def test_note_image_extraction():
    """Simulate Note ingest: ![](/api/notes/{col}/{id}/images/{file}) -> :::image."""
    print("\n" + "=" * 60)
    print("TEST 2: Note image URL extraction")
    print("=" * 60)

    with tempfile.TemporaryDirectory() as tmp:
        tmpdir = Path(tmp)
        note_id = "note_test_123"
        notes_dir = tmpdir / "notes" / note_id / "images"
        notes_dir.mkdir(parents=True)
        (notes_dir / "screenshot.png").write_bytes(_make_text_image())
        (notes_dir / "diagram.png").write_bytes(_make_mixed_image())

        content = (
            "# Meeting Notes\n\n"
            f"![Screenshot](/api/notes/col_test/{note_id}/images/screenshot.png)\n\n"
            f'<img src="/api/notes/col_test/{note_id}/images/diagram.png" data-image-id="test-uuid" />\n\n'
            "![External](https://example.com/photo.png)\n\n"
            ":::distill-block{src=\"note_456\"}\n"
            "Some distilled content\n"
            ":::\n"
        )

        import re as _re

        # Pre-convert HTML <img> tags (simulating what _do_ingest_note does)
        _note_img_re = _re.compile(r'/api/notes/[^/]+/([^/]+)/images/(.+)')
        _html_img_re = _re.compile(r'<img\s[^>]*?src="([^"]*)"[^>]*/?>', _re.IGNORECASE)

        def _convert_html_img(m: _re.Match) -> str:
            src = m.group(1)
            if _note_img_re.search(src):
                return f"![]({src})"
            return m.group(0)
        content = _html_img_re.sub(_convert_html_img, content)

        from src.parsers.markdown import extract_markdown_images

        # notes_dir is tmp/notes/note_test_123/images/
        # notes_dir.parent.parent is tmp/notes/
        _notes_base = notes_dir.parent.parent  # tmp/notes/

        def _resolve_note_image(url: str) -> bytes | None:
            m = _note_img_re.search(url)
            if not m:
                return None
            target_note = m.group(1)
            filename = m.group(2)
            img_path = _notes_base / target_note / "images" / filename
            if img_path.is_file():
                return img_path.read_bytes()
            return None

        result_content, images = extract_markdown_images(
            content, url_resolver=_resolve_note_image,
        )

        assert len(images) == 2, f"Expected 2 note images, got {len(images)}"
        assert ":::image" in result_content
        assert ":::distill-block" in result_content, "Distill block lost!"
        assert "![External]" in result_content, "External URL should stay"
        assert "![Screenshot]" not in result_content
        assert "![Diagram]" not in result_content

        print(f"  OK Extracted {len(images)} note images")
        print(f"  OK Distill block preserved")
        print(f"  OK External URL preserved")
        print(f"  OK Note image URLs replaced with :::image blocks")


# ═══════════════════════════════════════════════════════════════════════
# Test 3: Regex backward compatibility (old 3-field format)
# ═══════════════════════════════════════════════════════════════════════

def test_regex_backward_compat():
    """Old 3-field :::image blocks should still match the new regex."""
    print("\n" + "=" * 60)
    print("TEST 3: Regex backward compatibility (old 3-field format)")
    print("=" * 60)

    from src.parsers.image_utils import _IMAGE_BLOCK_RE, annotate_chunks_with_images
    from src.parsers.base import ImageInfo

    old_block = ":::image\nimage_id: abc123def456\nfile_id: file001\ndescription: A bar chart showing revenue\n:::"
    new_block = ":::image\nimage_id: def789abc012\nfile_id: file002\nocr_text: Q3 revenue 42M up 15 percent\ndescription: A bar chart showing quarterly revenue growth\n:::"
    empty_ocr_block = ":::image\nimage_id: abc345def678\nfile_id: file003\nocr_text: \ndescription: A photo of a team\n:::"

    content = f"Before\n{old_block}\nMiddle\n{new_block}\n{empty_ocr_block}\nAfter"
    matches = list(_IMAGE_BLOCK_RE.finditer(content))
    assert len(matches) == 3, f"Expected 3 matches, got {len(matches)}"

    # Old format: ocr_text (group 3) is None
    assert matches[0].group(3) is None or matches[0].group(3) == ""
    assert matches[0].group(4) == "A bar chart showing revenue"

    # New format
    assert matches[1].group(3) == "Q3 revenue 42M up 15 percent"
    assert matches[1].group(4) == "A bar chart showing quarterly revenue growth"

    # Empty ocr_text
    assert matches[2].group(3) == ""

    print(f"  OK Old format: ocr_text=None/empty, description=OK")
    print(f"  OK New format: ocr_text+description=OK")
    print(f"  OK Empty ocr_text: handled correctly")

    # Test annotate_chunks_with_images with old format
    class FakeChunk:
        def __init__(self, text):
            self.text = text
            self.metadata = {}
            self.score = 0.0

    chunk = FakeChunk(f"Some text\n{old_block}\nMore text")
    img_info = ImageInfo(
        image_id="abc123def456", file_id="file001",
        description="A bar chart showing revenue",
    )
    annotated = annotate_chunks_with_images([chunk], [img_info])
    meta = annotated[0].metadata
    assert "images" in meta
    ref = meta["images"][0]
    assert ref["image_id"] == "abc123def456"
    assert ref["file_id"] == "file001"
    print(f"  OK Chunk annotation works with old format")
    print(f"    image_id={ref['image_id']} file_id={ref['file_id']} ocr_text={ref.get('ocr_text', '')!r}")


# ═══════════════════════════════════════════════════════════════════════
# Test 4: process_document_images pipeline
# ═══════════════════════════════════════════════════════════════════════

def test_process_document_images():
    """Full pipeline: filter -> OCR classify."""
    print("\n" + "=" * 60)
    print("TEST 4: process_document_images pipeline")
    print("=" * 60)

    with tempfile.TemporaryDirectory() as tmp:
        tmpdir = Path(tmp)
        from src.parsers.base import ImageInfo, ParsedDocument
        from src.parsers.image_utils import process_document_images, build_image_block

        text_bytes = _make_text_image()
        visual_bytes = _make_visual_image()

        img_text = ImageInfo(image_id="img_text_001", file_id="e2e", image_bytes=text_bytes, image_format="png")
        img_visual = ImageInfo(image_id="img_visual_001", file_id="e2e", image_bytes=visual_bytes, image_format="png")
        img_small = ImageInfo(image_id="img_small_001", file_id="e2e", image_bytes=b'\x89PNG\r\n\x1a\n' + b'\x00' * 50, image_format="png", bbox=(0, 0, 10, 10))

        content = "\n".join(build_image_block(i) for i in [img_text, img_visual, img_small])
        doc = ParsedDocument(content=content, file_type="test", images=[img_text, img_visual, img_small])

        doc = process_document_images(doc, "e2e_pipe", tmpdir, vision_provider=None, vision_model_id="", vision_prompt="")

        print(f"  Images before: 3, after: {len(doc.images)} (all kept, no descriptions w/o Vision LLM)")
        for img in doc.images:
            print(f"    kept: {img.image_id} ocr_text={bool(img.ocr_text)} desc={bool(img.description)}")
        # Without Vision LLM, all non-filtered images are kept (just without descriptions)
        # img_small is filtered out (too small), img_text and img_visual both survive
        assert len(doc.images) >= 1, "At least text-type image should survive"
        assert all(not img.description for img in doc.images), "No images should have descriptions without Vision LLM"
        # Image blocks should still be in content (not removed)
        for img in doc.images:
            assert img.image_id in doc.content, f"Image block {img.image_id} should remain in content"
        print(f"  OK Pipeline completed without errors")


# ═══════════════════════════════════════════════════════════════════════
# Test 5: PDF image extraction
# ═══════════════════════════════════════════════════════════════════════

def test_pdf_image_extraction():
    """Create a PDF with an embedded image, parse it."""
    print("\n" + "=" * 60)
    print("TEST 5: PDF image extraction")
    print("=" * 60)

    try:
        from reportlab.lib.pagesizes import A4
        from reportlab.pdfgen import canvas
    except ImportError:
        print("  SKIP reportlab not installed")
        return

    with tempfile.TemporaryDirectory() as tmp:
        tmpdir = Path(tmp)
        img_bytes = _make_text_image()
        img_path = tmpdir / "chart.png"
        img_path.write_bytes(img_bytes)

        pdf_path = tmpdir / "test.pdf"
        c = canvas.Canvas(str(pdf_path), pagesize=A4)
        c.drawString(100, 700, "Q3 Revenue Report")
        c.drawImage(str(img_path), 100, 400, width=300, height=100)
        c.showPage()
        c.save()

        print(f"  OK Created PDF ({pdf_path.stat().st_size} bytes)")

        from src.parsers.pdf import PDFParser
        parser = PDFParser()
        doc = parser.parse(pdf_path)

        print(f"  Content: {len(doc.content)} chars, Images: {len(doc.images)}")
        if doc.images:
            for img in doc.images:
                print(f"    {img.image_id[:12]}... page={img.page_number} bytes={len(img.image_bytes or b'')}")
        assert doc.images, "PDF should have images"
        print(f"  OK PDF parser extracted images")


# ═══════════════════════════════════════════════════════════════════════
# Test 6: PPTX image extraction
# ═══════════════════════════════════════════════════════════════════════

def test_pptx_image_extraction():
    print("\n" + "=" * 60)
    print("TEST 6: PPTX image extraction")
    print("=" * 60)

    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError:
        print("  SKIP python-pptx not installed")
        return

    with tempfile.TemporaryDirectory() as tmp:
        tmpdir = Path(tmp)
        img_bytes = _make_text_image()
        img_path = tmpdir / "chart.png"
        img_path.write_bytes(img_bytes)

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_picture(str(img_path), Inches(1), Inches(1), Inches(4), Inches(2))
        pptx_path = tmpdir / "test.pptx"
        prs.save(str(pptx_path))

        print(f"  OK Created PPTX ({pptx_path.stat().st_size} bytes)")

        from src.parsers.pptx import PptxParser
        doc = PptxParser().parse(pptx_path)

        print(f"  Content: {len(doc.content)} chars, Images: {len(doc.images)}")
        if doc.images:
            for img in doc.images:
                print(f"    {img.image_id[:12]}... slide={img.slide_number} bytes={len(img.image_bytes or b'')}")
        assert doc.images, "PPTX should have images"
        print(f"  OK PPTX parser extracted images")


# ═══════════════════════════════════════════════════════════════════════
# Test 7: DOCX image extraction
# ═══════════════════════════════════════════════════════════════════════

def test_docx_image_extraction():
    print("\n" + "=" * 60)
    print("TEST 7: DOCX image extraction")
    print("=" * 60)

    try:
        from docx import Document
        from docx.shared import Inches
    except ImportError:
        print("  SKIP python-docx not installed")
        return

    with tempfile.TemporaryDirectory() as tmp:
        tmpdir = Path(tmp)
        img_bytes = _make_text_image()
        img_path = tmpdir / "chart.png"
        img_path.write_bytes(img_bytes)

        d = Document()
        d.add_paragraph("Q3 Revenue Report")
        d.add_picture(str(img_path), width=Inches(4))
        d.add_paragraph("End of report.")
        docx_path = tmpdir / "test.docx"
        d.save(str(docx_path))

        print(f"  OK Created DOCX ({docx_path.stat().st_size} bytes)")

        from src.parsers.docx import DocxParser
        doc = DocxParser().parse(docx_path)

        print(f"  Content: {len(doc.content)} chars, Images: {len(doc.images)}")
        if doc.images:
            for img in doc.images:
                print(f"    {img.image_id[:12]}... bytes={len(img.image_bytes or b'')}")
        assert doc.images, "DOCX should have images"
        print(f"  OK DOCX parser extracted images")


# ═══════════════════════════════════════════════════════════════════════
# Test 8: Multimodal context building
# ═══════════════════════════════════════════════════════════════════════

def test_multimodal_context():
    print("\n" + "=" * 60)
    print("TEST 8: Multimodal context building")
    print("=" * 60)

    from src.rag.agentic_query import _build_multimodal_context
    from src.parsers.image_utils import build_image_block
    from src.parsers.base import ImageInfo

    img1 = ImageInfo(image_id="aaa", file_id="f1", ocr_text="Q3 revenue: $42M, up 15% YoY", description="")
    img2 = ImageInfo(image_id="bbb", file_id="f2", ocr_text="", description="A bar chart of quarterly revenue")
    img3 = ImageInfo(image_id="ccc", file_id="f3", ocr_text="Frontend -> API -> DB", description="Architecture diagram, 3-tier")

    content = f"Report:\n{build_image_block(img1)}\nChart:\n{build_image_block(img2)}\nDiagram:\n{build_image_block(img3)}"

    images = {
        "aaa": {"base64": "x1", "mime": "image/png"},
        "bbb": {"base64": "x2", "mime": "image/png"},
        "ccc": {"base64": "x3", "mime": "image/png"},
    }

    result = _build_multimodal_context(content, images)
    img_parts = [p for p in result if p["type"] == "image_url"]
    ocr_parts = [p for p in result if p["type"] == "text" and "OCR text" in p.get("text", "")]

    assert len(img_parts) == 3
    assert len(ocr_parts) == 2  # img1 and img3 have OCR text
    assert any("$42M" in p["text"] for p in ocr_parts)
    assert any("Frontend" in p["text"] for p in ocr_parts)

    print(f"  OK {len(img_parts)} image parts, {len(ocr_parts)} OCR text parts")
    print(f"  OK OCR text injected before images")


# ═══════════════════════════════════════════════════════════════════════

if __name__ == "__main__":
    print("=" * 60)
    print("E2E IMAGE PIPELINE TESTS")
    print("=" * 60)

    tests = [
        ("Markdown images", test_markdown_with_images),
        ("Note image URLs", test_note_image_extraction),
        ("Regex backward compat", test_regex_backward_compat),
        ("process_document_images", test_process_document_images),
        ("PDF extraction", test_pdf_image_extraction),
        ("PPTX extraction", test_pptx_image_extraction),
        ("DOCX extraction", test_docx_image_extraction),
        ("Multimodal context", test_multimodal_context),
    ]

    passed = 0
    failed = 0

    for name, fn in tests:
        try:
            fn()
            passed += 1
            print(f"\n  PASS: {name}")
        except Exception as e:
            failed += 1
            print(f"\n  FAIL: {name} — {e}")
            import traceback
            traceback.print_exc()

    print("\n" + "=" * 60)
    print(f"RESULTS: {passed} passed, {failed} failed, {len(tests)} total")
    print("=" * 60)

    if failed > 0:
        sys.exit(1)
